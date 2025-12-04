"""
Model Context Protocol (MCP) Server Implementation for AI Agent

The MCP Server acts as a standardized bridge between AI models and external systems,
providing unified connectivity, real-time access, and scalable integration.

Key Features:
- Unified Connectivity: Single standardized interface for all external tools and data sources
- Real-Time Access: Enables AI models to pull fresh, live information instead of static training data
- Scalability: Minimal setup required, reducing deployment time and complexity
- Flexibility: Extensible architecture that can integrate with any external service

Architecture:
The MCP Server implements the Model Context Protocol, allowing AI models to:
1. Discover available tools through standardized tool descriptions
2. Execute tools with structured parameters
3. Receive real-time results formatted for AI consumption
4. Access multiple data sources (files, web, code analysis) through one interface

This implementation provides tools for:
- File operations (read, write, list, search)
- Code analysis and search
- Web search with caching and optimization
- Command execution
- Directory tree navigation

The server follows the MCP protocol standard, ensuring compatibility with any
MCP-compliant AI model or client.
"""

import asyncio
import json
import os
import re
import subprocess
import time
import logging
import aiofiles
import aiohttp
from typing import Any, Dict, List, Optional, Tuple
from pathlib import Path
from datetime import datetime
from urllib.parse import urlparse

logger = logging.getLogger(__name__)

# Import here to avoid circular dependencies
MCP_AVAILABLE = False
try:
    from mcp.server import Server
    from mcp.server.stdio import stdio_server
    from mcp.types import Tool, TextContent
    MCP_AVAILABLE = True
except ImportError:
    MCP_AVAILABLE = False
    # Create minimal stubs for when MCP is not available
    class Server:
        pass
    class stdio_server:
        pass
    
    # Simple Tool and TextContent classes for when MCP SDK is not available
    class Tool:
        def __init__(self, name: str, description: str, inputSchema: Dict[str, Any]):
            self.name = name
            self.description = description
            self.inputSchema = inputSchema
    
    class TextContent:
        def __init__(self, type: str, text: str):
            self.type = type
            self.text = text


class MCPServerTools:
    """
    MCP Server Tools - Bridge between AI models and external systems
    
    This class implements the Model Context Protocol server, providing a unified
    interface for AI models to access external tools and data sources. It acts as
    the backend system that handles requests from AI models and routes them to
    the appropriate tools, APIs, or services.
    
    The MCP server provides:
    - Standardized tool discovery and execution
    - Real-time data access (files, web, code)
    - Caching for performance optimization
    - Error handling and validation
    
    This follows the MCP protocol standard, ensuring seamless integration with
    any MCP-compliant AI model or client.
    """
    
    def __init__(self, file_service=None, code_analyzer=None, web_search_enabled=True, web_search_service=None, workspace_root=None, location_service=None, memory_service=None):
        self.file_service = file_service
        self.code_analyzer = code_analyzer
        self.web_search_enabled = web_search_enabled
        self.workspace_root = os.path.abspath(workspace_root) if workspace_root else os.getcwd()
        self._web_search_service = web_search_service  # Shared web search service instance
        self._location_service = location_service  # Location service for weather/news
        self._memory_service = memory_service  # Memory service for saving memories
        self._current_images = []  # Store current images for identify_image tool access
        try:
            self._cache_ttl_seconds = max(1, int(os.getenv("MCP_CACHE_TTL_SECONDS", "4")))
        except ValueError:
            self._cache_ttl_seconds = 4
        self._dir_cache: Dict[str, Dict[str, Any]] = {}
        self._tree_cache: Dict[str, Dict[str, Any]] = {}
        
        # Command execution loop prevention
        # Track recent commands to prevent infinite loops
        self._command_history: List[Dict[str, Any]] = []  # List of {command, timestamp} dicts
        self._max_command_history = 20  # Keep last 20 commands
        self._loop_detection_window = 30  # Seconds to look back for loops
        self._max_repeats = 3  # Maximum times same command can run in window
        
        # Tool metadata
        self.server_version = "1.0.0"
        self.server_capabilities = {
            "caching": True,
            "validation": True,
            "analytics": True,
            "error_recovery": True
        }

    def _build_cache_key(self, path: str, suffix: str = "") -> str:
        normalized = os.path.abspath(path)
        return f"{normalized}::{suffix}" if suffix else normalized

    def _get_cached_text(self, cache: Dict[str, Dict[str, Any]], key: str) -> Optional[str]:
        entry = cache.get(key)
        if not entry:
            return None
        if (time.time() - entry.get("ts", 0)) > self._cache_ttl_seconds:
            cache.pop(key, None)
            return None
        return entry.get("text")

    def _set_cached_text(self, cache: Dict[str, Dict[str, Any]], key: str, text: str) -> None:
        cache[key] = {"ts": time.time(), "text": text}

    def _invalidate_structure_caches(self) -> None:
        self._dir_cache.clear()
        self._tree_cache.clear()
    
    def _check_command_loop(self, command: str) -> Tuple[bool, Optional[str]]:
        """
        Check if a command would create a loop by running the same command too many times.
        
        Returns:
            Tuple of (is_loop, error_message)
            - is_loop: True if this command would create a loop
            - error_message: Error message if loop detected, None otherwise
        """
        current_time = time.time()
        normalized_command = command.strip()
        
        # Clean old entries outside the detection window
        self._command_history = [
            entry for entry in self._command_history
            if (current_time - entry.get("timestamp", 0)) <= self._loop_detection_window
        ]
        
        # Count how many times this exact command was run recently
        recent_count = sum(
            1 for entry in self._command_history
            if entry.get("command", "").strip() == normalized_command
        )
        
        # Check if this would exceed the maximum repeats
        if recent_count >= self._max_repeats:
            return True, (
                f"LOOP PREVENTION: Command '{normalized_command}' has been executed "
                f"{recent_count} times in the last {self._loop_detection_window} seconds. "
                f"Maximum allowed repeats: {self._max_repeats}. "
                f"This prevents infinite loops. Please wait a moment or modify the command."
            )
        
        return False, None
    
    def _record_command_execution(self, command: str) -> None:
        """Record a command execution in the history for loop detection."""
        current_time = time.time()
        normalized_command = command.strip()
        
        # Add to history
        self._command_history.append({
            "command": normalized_command,
            "timestamp": current_time
        })
        
        # Keep only the most recent entries
        if len(self._command_history) > self._max_command_history:
            self._command_history = self._command_history[-self._max_command_history:]
    
    def set_workspace_root(self, workspace_path: str) -> None:
        """
        Update the workspace root path for MCP tool operations.
        This ensures tools operate within the correct workspace directory.
        """
        if workspace_path and workspace_path.strip():
            # Normalize the path - handle both absolute and relative paths
            normalized = workspace_path.strip().replace('\\', '/')
            # If it's a relative path, resolve it relative to current working directory
            if not os.path.isabs(normalized):
                resolved_path = os.path.abspath(os.path.join(os.getcwd(), normalized))
            else:
                resolved_path = os.path.abspath(normalized)
            
            # Verify the path exists and is a directory
            if os.path.exists(resolved_path) and os.path.isdir(resolved_path):
                old_root = self.workspace_root
                self.workspace_root = resolved_path
                # Invalidate caches when workspace changes
                if old_root != self.workspace_root:
                    self._invalidate_structure_caches()
                    logger.info(f"[MCP] Workspace root updated: {old_root} -> {self.workspace_root}")
            else:
                logger.warning(f"[MCP] Workspace path does not exist or is not a directory: {workspace_path} (resolved: {resolved_path}), keeping current: {self.workspace_root}")
        else:
            # Reset to default (current working directory) if no path provided
            old_root = self.workspace_root
            self.workspace_root = os.getcwd()
            if old_root != self.workspace_root:
                self._invalidate_structure_caches()
                logger.info(f"[MCP] Workspace root reset to default: {self.workspace_root}")
    
    def get_tools(self) -> List[Tool]:
        """
        Get list of available MCP tools
        
        This method implements the MCP protocol's tool discovery mechanism.
        It returns a standardized list of tools that AI models can use, with
        complete schema information for each tool.
        
        Returns:
            List of Tool objects following the MCP protocol specification
        """
        # Tool class is always available - either from mcp.types or from the fallback stub
        # defined at module level (lines 59-63). No need to check MCP_AVAILABLE or import.
        
        tools = [
            Tool(
                name="read_file",
                description="Read the contents of a file. Path can be relative to workspace root or absolute.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path to the file to read (relative to workspace or absolute)"
                        }
                    },
                    "required": ["path"]
                }
            ),
            Tool(
                name="write_file",
                description="Write content to a file. Creates the file if it doesn't exist. In ASK mode, this tool is disabled.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path to the file to write (relative to workspace or absolute)"
                        },
                        "content": {
                            "type": "string",
                            "description": "Content to write to the file"
                        }
                    },
                    "required": ["path", "content"]
                }
            ),
            Tool(
                name="list_directory",
                description="List files and directories in a given path",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path to list (defaults to workspace root)",
                            "default": "."
                        }
                    }
                }
            ),
            Tool(
                name="search_files",
                description="Search for files by name pattern in the workspace",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "query": {
                            "type": "string",
                            "description": "File name pattern to search for (supports wildcards)"
                        },
                        "path": {
                            "type": "string",
                            "description": "Directory to search in (defaults to workspace root)",
                            "default": "."
                        }
                    },
                    "required": ["query"]
                }
            ),
            Tool(
                name="get_file_tree",
                description="Get the directory tree structure of the workspace or a specific path",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path to get tree for (defaults to workspace root)",
                            "default": "."
                        },
                        "max_depth": {
                            "type": "integer",
                            "description": "Maximum depth to traverse (default: 4)",
                            "default": 4
                        }
                    }
                }
            ),
            Tool(
                name="analyze_code",
                description="Analyze code in a file: extract functions, classes, imports, and dependencies",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path to the code file to analyze"
                        }
                    },
                    "required": ["path"]
                }
            ),
            Tool(
                name="grep_code",
                description="Search for text patterns in code files across the workspace",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "pattern": {
                            "type": "string",
                            "description": "Text pattern or regex to search for"
                        },
                        "path": {
                            "type": "string",
                            "description": "Directory to search in (defaults to workspace root)",
                            "default": "."
                        },
                        "file_extensions": {
                            "type": "array",
                            "items": {"type": "string"},
                            "description": "File extensions to limit search to (e.g., ['.py', '.js'])"
                        }
                    },
                    "required": ["pattern"]
                }
            ),
            Tool(
                name="execute_command",
                description="Execute a shell command in the workspace directory and return the output. Useful for investigating issues, testing code temporarily, checking file contents, running tests, or debugging. The command output (stdout and stderr) will be visible in the chat. Loop prevention is enabled to prevent infinite command execution.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "command": {
                            "type": "string",
                            "description": "Shell command to execute (e.g., 'ls -la', 'python test.py', 'npm test', 'git status')"
                        },
                        "timeout": {
                            "type": "integer",
                            "description": "Timeout in seconds (default: 30)",
                            "default": 30
                        }
                    },
                    "required": ["command"]
                }
            ),
            Tool(
                name="download_file",
                description="Download a file from a URL and save it to the workspace. Supports HTTP and HTTPS URLs.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "url": {
                            "type": "string",
                            "description": "URL of the file to download (must be HTTP or HTTPS)"
                        },
                        "path": {
                            "type": "string",
                            "description": "Path where to save the file (relative to workspace or absolute). If not provided, filename will be extracted from URL."
                        },
                        "timeout": {
                            "type": "integer",
                            "description": "Download timeout in seconds (default: 60)",
                            "default": 60
                        }
                    },
                    "required": ["url"]
                }
            ),
            Tool(
                name="create_document",
                description="Create a Microsoft Word-like document (.docx) with formatted text, headings, paragraphs, lists, and tables. Similar to Google Docs or Microsoft Word.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path where to save the document (relative to workspace or absolute). Should end with .docx"
                        },
                        "title": {
                            "type": "string",
                            "description": "Document title"
                        },
                        "content": {
                            "type": "string",
                            "description": "Document content. Can include markdown-style formatting: # for headings, * for lists, ** for bold, etc."
                        },
                        "author": {
                            "type": "string",
                            "description": "Document author (optional)"
                        }
                    },
                    "required": ["path", "content"]
                }
            ),
            Tool(
                name="create_slide",
                description="Create a single PowerPoint slide (.pptx) with enhanced formatting and design. Supports bullet points (use '- ' or '* '), headings (use '# ' or '## '), bold text (use **text**), nested bullets, and proper text styling. Content is automatically formatted with professional styling, colors, and spacing.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path where to save the slide presentation (relative to workspace or absolute). Should end with .pptx"
                        },
                        "title": {
                            "type": "string",
                            "description": "Slide title (will be styled with large, bold, blue text)"
                        },
                        "content": {
                            "type": "string",
                            "description": "Slide content with formatting: Use '- ' or '* ' for bullet points, '## ' for subheadings, '# ' for main headings, '**text**' for bold text. Content is automatically formatted with proper fonts, colors, and spacing."
                        },
                        "layout": {
                            "type": "string",
                            "description": "Slide layout: 'title_only', 'title_content', 'blank', 'title_slide' (default: 'title_content')",
                            "enum": ["title_only", "title_content", "blank", "title_slide"],
                            "default": "title_content"
                        }
                    },
                    "required": ["path", "title"]
                }
            ),
            Tool(
                name="create_presentation",
                description="Create a full PowerPoint presentation (.pptx) with multiple slides and enhanced formatting. Each slide supports bullet points (use '- ' or '* '), headings (use '# ' or '## '), bold text (use **text**), nested bullets, and professional styling. All slides are automatically formatted with proper fonts, colors, spacing, and visual hierarchy.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "path": {
                            "type": "string",
                            "description": "Path where to save the presentation (relative to workspace or absolute). Should end with .pptx"
                        },
                        "title": {
                            "type": "string",
                            "description": "Presentation title (shown on title slide with large, bold, centered text)"
                        },
                        "slides": {
                            "type": "array",
                            "description": "Array of slide objects, each with 'title' and 'content' fields. Content supports formatting: Use '- ' or '* ' for bullet points, '## ' for subheadings, '# ' for main headings, '**text**' for bold text.",
                            "items": {
                                "type": "object",
                                "properties": {
                                    "title": {"type": "string", "description": "Slide title (styled with large, bold, blue text)"},
                                    "content": {"type": "string", "description": "Slide content with formatting support (bullet points, headings, bold text)"},
                                    "layout": {
                                        "type": "string",
                                        "enum": ["title_only", "title_content", "blank", "title_slide"],
                                        "default": "title_content"
                                    }
                                },
                                "required": ["title"]
                            }
                        },
                        "author": {
                            "type": "string",
                            "description": "Presentation author (optional, shown on title slide)"
                        }
                    },
                    "required": ["path", "title", "slides"]
                }
            ),
            Tool(
                name="predict_price",
                description="Predict future price movements for cryptocurrency or forex assets using advanced technical analysis, trend analysis, and statistical forecasting. Returns price predictions with confidence levels, trend analysis, and risk assessment.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "asset": {
                            "type": "string",
                            "description": "Asset identifier (e.g., 'bitcoin', 'btc', 'ethereum', 'eth', 'eur/usd', 'gbp/usd'). Supports crypto and forex pairs."
                        },
                        "asset_type": {
                            "type": "string",
                            "description": "Asset type: 'crypto' or 'forex' (auto-detected if not provided)",
                            "enum": ["crypto", "forex"]
                        },
                        "days_ahead": {
                            "type": "integer",
                            "description": "Number of days to predict ahead (1-30, default: 7)",
                            "default": 7,
                            "minimum": 1,
                            "maximum": 30
                        },
                        "include_analysis": {
                            "type": "boolean",
                            "description": "Include detailed technical analysis (default: true)",
                            "default": True
                        }
                    },
                    "required": ["asset"]
                }
            ),
            Tool(
                name="identify_image",
                description="Identify and describe the contents of an image using AI vision capabilities. Analyzes images to detect objects, text, UI elements, code, diagrams, screenshots, and other visual content. Use this tool when users upload images and ask 'what is in this image' or similar questions. If images are attached to the current message, they are automatically available - you can call this tool without providing image_data, or provide image_data explicitly.",
                inputSchema={
                    "type": "object",
                    "properties": {
                        "image_data": {
                            "type": "string",
                            "description": "Base64-encoded image data (data URL format: data:image/...;base64,...) or base64 string without prefix. Optional if images are attached to the current message (they will be used automatically)."
                        },
                        "query": {
                            "type": "string",
                            "description": "Optional specific question about the image (e.g., 'what code is shown?', 'identify the UI elements', 'what is this diagram about?', 'what is in this image?')"
                        }
                    },
                    "required": []
                }
            ),
        ]
        
        if self.web_search_enabled:
            tools.append(
                Tool(
                    name="web_search",
                    description="Search the web using DuckDuckGo with enhanced features: result caching, relevance scoring, and query optimization. Returns search results with titles, URLs, and snippets.",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "query": {
                                "type": "string",
                                "description": "Search query (will be optimized automatically for better results)"
                            },
                            "max_results": {
                                "type": "integer",
                                "description": "Maximum number of results to return (default: 5, max: 20)",
                                "default": 5,
                                "minimum": 1,
                                "maximum": 20
                            },
                            "search_type": {
                                "type": "string",
                                "description": "Type of search: 'text' (default), 'news', or 'images'",
                                "enum": ["text", "news", "images"],
                                "default": "text"
                            }
                        },
                        "required": ["query"]
                    }
                )
            )
        
        # Add location-based tools if location service is available
        if self._location_service:
            tools.extend([
                Tool(
                    name="get_user_location",
                    description="Detect the user's current location based on IP address geolocation. Returns city, country, coordinates, and timezone information.",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "ip_address": {
                                "type": "string",
                                "description": "Optional IP address. If not provided, uses the current connection's IP address."
                            }
                        }
                    }
                ),
                Tool(
                    name="get_weather",
                    description="Get current weather information for a location. Can use city name or coordinates. Returns temperature, conditions, humidity, wind speed, and more.",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "city": {
                                "type": "string",
                                "description": "City name (e.g., 'London', 'New York'). Either city or coordinates must be provided."
                            },
                            "latitude": {
                                "type": "number",
                                "description": "Latitude coordinate. Required if city is not provided."
                            },
                            "longitude": {
                                "type": "number",
                                "description": "Longitude coordinate. Required if city is not provided."
                            },
                            "units": {
                                "type": "string",
                                "description": "Temperature units: 'metric' for Celsius (default), 'imperial' for Fahrenheit",
                                "enum": ["metric", "imperial"],
                                "default": "metric"
                            }
                        }
                    }
                ),
                Tool(
                    name="get_news",
                    description="Get news articles for a location. Can filter by city, country, or search query. Returns recent news articles with titles, descriptions, and URLs.",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "city": {
                                "type": "string",
                                "description": "City name to get news for (e.g., 'London', 'New York')"
                            },
                            "country": {
                                "type": "string",
                                "description": "Country name or code (e.g., 'US', 'GB', 'United States')"
                            },
                            "query": {
                                "type": "string",
                                "description": "Optional search query to filter news articles"
                            },
                            "max_results": {
                                "type": "integer",
                                "description": "Maximum number of news articles to return (default: 10, max: 50)",
                                "default": 10,
                                "minimum": 1,
                                "maximum": 50
                            }
                        }
                    }
                ),
            ])
        
        # Add memory tools if memory service is available
        if self._memory_service:
            tools.extend([
                Tool(
                    name="save_memory",
                    description="Save a memory that the user wants you to remember. Use this when the user explicitly asks you to remember, save, or keep in mind something. Examples: 'Remember that my name is John', 'Save that I prefer dark mode', 'Keep in mind that I'm a Python developer'.",
                    inputSchema={
                        "type": "object",
                        "properties": {
                            "content": {
                                "type": "string",
                                "description": "The memory content to save. Should be a clear, concise statement of what to remember (e.g., 'User's name is John', 'User prefers dark mode', 'User is a Python developer')."
                            }
                        },
                        "required": ["content"]
                    }
                ),
            ])
        
        return tools
    
    def get_server_info(self) -> Dict[str, Any]:
        """
        Get MCP server information and capabilities
        
        Returns:
            Dictionary containing server metadata including version, capabilities,
            and available tools count
        """
        tools = self.get_tools()
        return {
            "version": self.server_version,
            "capabilities": self.server_capabilities,
            "tool_count": len(tools),
            "tools": [tool.name for tool in tools],
            "workspace_root": self.workspace_root,
            "web_search_enabled": self.web_search_enabled
        }
    
    async def execute_tool(self, tool_name: str, arguments: Dict[str, Any], allow_write: bool = True) -> List[TextContent]:
        """
        Execute an MCP tool and return results
        
        This is the core method that implements the MCP protocol's tool execution.
        It routes tool calls from AI models to the appropriate backend services,
        providing real-time access to external systems.
        
        The method:
        1. Validates the tool name and arguments
        2. Routes to the appropriate tool handler
        3. Executes the tool with proper error handling
        4. Returns standardized results in MCP TextContent format
        
        This unified interface allows AI models to access multiple external
        systems (files, web, code analysis) without needing custom integrations
        for each service.
        
        Args:
            tool_name: Name of the tool to execute (must be in get_tools())
            arguments: Tool-specific parameters
            allow_write: Whether write operations are permitted (for ASK mode)
            
        Returns:
            List of TextContent objects with tool execution results
        """
        if not MCP_AVAILABLE:
            logger.error("MCP SDK not available")
            return [TextContent(
                type="text",
                text=f"MCP SDK not available. Please install with: pip install mcp"
            )]
        
        execution_start = time.time()
        logger.info(f"Executing MCP tool: {tool_name} with arguments: {arguments}")
        
        try:
            if tool_name == "read_file":
                return await self._read_file(arguments.get("path", ""))
            elif tool_name == "write_file":
                if not allow_write:
                    return [TextContent(
                        type="text",
                        text="ERROR: Write operations are disabled in ASK mode. Switch to Agent mode to modify files."
                    )]
                return await self._write_file(arguments.get("path", ""), arguments.get("content", ""))
            elif tool_name == "list_directory":
                return await self._list_directory(arguments.get("path", "."))
            elif tool_name == "search_files":
                return await self._search_files(arguments.get("query", ""), arguments.get("path", "."))
            elif tool_name == "get_file_tree":
                return await self._get_file_tree(
                    arguments.get("path", "."),
                    arguments.get("max_depth", 4)
                )
            elif tool_name == "analyze_code":
                return await self._analyze_code(arguments.get("path", ""))
            elif tool_name == "grep_code":
                return await self._grep_code(
                    arguments.get("pattern", ""),
                    arguments.get("path", "."),
                    arguments.get("file_extensions", [])
                )
            elif tool_name == "execute_command":
                if not allow_write:
                    return [TextContent(
                        type="text",
                        text="ERROR: Command execution is disabled in ASK mode. Switch to Agent mode to run commands."
                    )]
                return await self._execute_command(
                    arguments.get("command", ""),
                    arguments.get("timeout", 30)
                )
            elif tool_name == "web_search":
                return await self._web_search(
                    arguments.get("query", ""),
                    arguments.get("max_results", 5),
                    arguments.get("search_type", "text")
                )
            elif tool_name == "download_file":
                if not allow_write:
                    return [TextContent(
                        type="text",
                        text="ERROR: File downloads are disabled in ASK mode. Switch to Agent mode to download files."
                    )]
                return await self._download_file(
                    arguments.get("url", ""),
                    arguments.get("path"),
                    arguments.get("timeout", 60)
                )
            elif tool_name == "create_document":
                if not allow_write:
                    return [TextContent(
                        type="text",
                        text="ERROR: Document creation is disabled in ASK mode. Switch to Agent mode to create documents."
                    )]
                return await self._create_document(
                    arguments.get("path", ""),
                    arguments.get("content", ""),
                    arguments.get("title", ""),
                    arguments.get("author")
                )
            elif tool_name == "create_slide":
                if not allow_write:
                    return [TextContent(
                        type="text",
                        text="ERROR: Slide creation is disabled in ASK mode. Switch to Agent mode to create slides."
                    )]
                return await self._create_slide(
                    arguments.get("path", ""),
                    arguments.get("title", ""),
                    arguments.get("content", ""),
                    arguments.get("layout", "title_content")
                )
            elif tool_name == "create_presentation":
                if not allow_write:
                    return [TextContent(
                        type="text",
                        text="ERROR: Presentation creation is disabled in ASK mode. Switch to Agent mode to create presentations."
                    )]
                return await self._create_presentation(
                    arguments.get("path", ""),
                    arguments.get("title", ""),
                    arguments.get("slides", []),
                    arguments.get("author")
                )
            elif tool_name == "get_user_location":
                return await self._get_user_location(arguments.get("ip_address"))
            elif tool_name == "get_weather":
                return await self._get_weather(
                    arguments.get("city"),
                    arguments.get("latitude"),
                    arguments.get("longitude"),
                    arguments.get("units", "metric")
                )
            elif tool_name == "get_news":
                return await self._get_news(
                    arguments.get("city"),
                    arguments.get("country"),
                    arguments.get("query"),
                    arguments.get("max_results", 10)
                )
            elif tool_name == "save_memory":
                return await self._save_memory(arguments.get("content", ""))
            elif tool_name == "predict_price":
                return await self._predict_price(
                    arguments.get("asset", ""),
                    arguments.get("asset_type"),
                    arguments.get("days_ahead", 7),
                    arguments.get("include_analysis", True)
                )
            elif tool_name == "identify_image":
                # If no image_data provided, try to use current images from context
                image_data = arguments.get("image_data", "")
                if not image_data and self._current_images:
                    # Use first available image if none specified
                    image_data = self._current_images[0]
                    logger.info(f"[MCP] Using current image from context for identify_image")
                return await self._identify_image(
                    image_data,
                    arguments.get("query")
                )
            else:
                execution_time = time.time() - execution_start
                logger.warning(f"Unknown tool requested: {tool_name} (execution time: {execution_time:.3f}s)")
                return [TextContent(
                    type="text",
                    text=f"Unknown tool: {tool_name}. Available tools: {', '.join([t.name for t in self.get_tools()])}"
                )]
        except Exception as e:
            execution_time = time.time() - execution_start
            error_msg = f"Error executing {tool_name}: {str(e)}"
            logger.error(f"{error_msg} (execution time: {execution_time:.3f}s)", exc_info=True)
            return [TextContent(
                type="text",
                text=error_msg
            )]
    
    async def _read_file(self, path: str) -> List[TextContent]:
        """Read a file"""
        if not self.file_service:
            return [TextContent(type="text", text="File service not available")]
        
        try:
            # Normalize path
            if not os.path.isabs(path):
                path = os.path.join(self.workspace_root, path)
            
            content = await self.file_service.read_file(path)
            return [TextContent(type="text", text=f"File: {path}\n\n{content}")]
        except FileNotFoundError:
            return [TextContent(type="text", text=f"File not found: {path}")]
        except Exception as e:
            return [TextContent(type="text", text=f"Error reading file: {str(e)}")]
    
    async def _write_file(self, path: str, content: str) -> List[TextContent]:
        """Write to a file"""
        if not self.file_service:
            return [TextContent(type="text", text="File service not available")]
        
        try:
            # Normalize path
            if not os.path.isabs(path):
                path = os.path.join(self.workspace_root, path)
            
            await self.file_service.write_file(path, content)
            self._invalidate_structure_caches()
            return [TextContent(type="text", text=f"Successfully wrote to {path}")]
        except Exception as e:
            return [TextContent(type="text", text=f"Error writing file: {str(e)}")]
    
    async def _list_directory(self, path: str) -> List[TextContent]:
        """List directory contents"""
        if not self.file_service:
            return [TextContent(type="text", text="File service not available")]
        
        try:
            # Normalize path
            if not os.path.isabs(path) or path == ".":
                path = os.path.join(self.workspace_root, path) if path != "." else self.workspace_root

            cache_key = self._build_cache_key(path, "list")
            cached_text = self._get_cached_text(self._dir_cache, cache_key)
            if cached_text:
                return [TextContent(type="text", text=cached_text)]

            files = await self.file_service.list_directory(path)
            file_list = []
            for file_info in files:
                file_type = "DIR" if file_info.is_directory else "FILE"
                size = f"{file_info.size} bytes" if not file_info.is_directory else ""
                file_list.append(f"{file_type:4s} {file_info.name:50s} {size}")
            
            result = f"Directory: {path}\n\n" + "\n".join(file_list)
            self._set_cached_text(self._dir_cache, cache_key, result)
            return [TextContent(type="text", text=result)]
        except Exception as e:
            return [TextContent(type="text", text=f"Error listing directory: {str(e)}")]
    
    async def _search_files(self, query: str, path: str) -> List[TextContent]:
        """Search for files"""
        if not self.file_service:
            return [TextContent(type="text", text="File service not available")]
        
        try:
            if not os.path.isabs(path) or path == ".":
                path = os.path.join(self.workspace_root, path) if path != "." else self.workspace_root
            
            results = await self.file_service.search_files(query, path)
            if not results:
                return [TextContent(type="text", text=f"No files matching '{query}' were found in {path}.")]

            formatted_results = []
            for entry in results[:20]:
                entry_path = entry.get("path") or entry.get("name") or ""
                info_parts = []
                size_value = entry.get("size")
                if isinstance(size_value, int):
                    info_parts.append(f"{size_value} bytes")
                if entry.get("modified_time"):
                    info_parts.append(f"modified {entry['modified_time']}")
                meta = f" ({'; '.join(info_parts)})" if info_parts else ""
                formatted_results.append(f"- {entry_path}{meta}")

            result = f"Search results for '{query}' in {path}:\n\n" + "\n".join(formatted_results)
            if len(results) > 20:
                result += f"\n... and {len(results) - 20} more results"
            return [TextContent(type="text", text=result)]
        except Exception as e:
            return [TextContent(type="text", text=f"Error searching files: {str(e)}")]
    
    async def _get_file_tree(self, path: str, max_depth: int) -> List[TextContent]:
        """Get file tree structure"""
        if not self.file_service:
            return [TextContent(type="text", text="File service not available")]
        
        try:
            if not os.path.isabs(path) or path == ".":
                path = os.path.join(self.workspace_root, path) if path != "." else self.workspace_root

            cache_key = self._build_cache_key(path, f"tree:{max_depth}")
            cached_text = self._get_cached_text(self._tree_cache, cache_key)
            if cached_text:
                return [TextContent(type="text", text=cached_text)]

            tree = await self.file_service.get_project_structure(path, max_depth=max_depth)
            
            def format_tree(node, indent=0):
                lines = []
                prefix = "  " * indent + ("└── " if indent > 0 else "")
                lines.append(f"{prefix}{node['name']}")
                for child in node.get("children", [])[:20]:  # Limit children
                    lines.extend(format_tree(child, indent + 1))
                return lines
            
            tree_lines = format_tree(tree)
            result = f"File tree for {path}:\n\n" + "\n".join(tree_lines)
            self._set_cached_text(self._tree_cache, cache_key, result)
            return [TextContent(type="text", text=result)]
        except Exception as e:
            return [TextContent(type="text", text=f"Error getting file tree: {str(e)}")]
    
    async def _analyze_code(self, path: str) -> List[TextContent]:
        """Analyze code in a file"""
        if not self.code_analyzer:
            return [TextContent(type="text", text="Code analyzer not available")]
        
        try:
            if not os.path.isabs(path):
                path = os.path.join(self.workspace_root, path)
            
            analysis = await self.code_analyzer.analyze_file(path)
            result = json.dumps(analysis, indent=2)
            return [TextContent(type="text", text=f"Code analysis for {path}:\n\n{result}")]
        except Exception as e:
            return [TextContent(type="text", text=f"Error analyzing code: {str(e)}")]
    
    async def _grep_code(self, pattern: str, path: str, file_extensions: List[str]) -> List[TextContent]:
        """Search for patterns in code"""
        if not self.code_analyzer:
            return [TextContent(type="text", text="Code analyzer not available")]
        
        try:
            if not os.path.isabs(path) or path == ".":
                path = os.path.join(self.workspace_root, path) if path != "." else self.workspace_root
            
            # Use code analyzer's grep functionality if available
            # Otherwise, do a simple recursive search
            matches = []
            for root, dirs, files in os.walk(path):
                # Skip hidden directories
                dirs[:] = [d for d in dirs if not d.startswith('.')]
                
                for file in files:
                    if file_extensions:
                        if not any(file.endswith(ext) for ext in file_extensions):
                            continue
                    
                    file_path = os.path.join(root, file)
                    try:
                        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                            for line_num, line in enumerate(f, 1):
                                if pattern.lower() in line.lower():
                                    matches.append(f"{file_path}:{line_num}: {line.strip()}")
                    except Exception:
                        continue
                    
                    if len(matches) >= 50:  # Limit results
                        break
                
                if len(matches) >= 50:
                    break
            
            result = f"Grep results for '{pattern}' in {path}:\n\n" + "\n".join(matches[:50])
            if len(matches) > 50:
                result += f"\n... and {len(matches) - 50} more matches"
            return [TextContent(type="text", text=result)]
        except Exception as e:
            return [TextContent(type="text", text=f"Error grepping code: {str(e)}")]
    
    async def _execute_command(self, command: str, timeout: int) -> List[TextContent]:
        """
        Execute a shell command with loop prevention.
        
        This method:
        1. Checks for command loops to prevent infinite execution
        2. Executes the command in the workspace directory
        3. Returns formatted output (stdout and stderr) that will be visible in chat
        4. Records the command execution for loop detection
        """
        try:
            # Normalize command for loop detection
            normalized_command = command.strip()
            
            if not normalized_command:
                return [TextContent(
                    type="text",
                    text="ERROR: Command cannot be empty"
                )]
            
            # Check for loops before executing
            is_loop, loop_error = self._check_command_loop(normalized_command)
            if is_loop:
                return [TextContent(
                    type="text",
                    text=loop_error or "LOOP PREVENTION: Command execution blocked to prevent infinite loops"
                )]
            
            # Record command execution (before running, so we track attempts)
            self._record_command_execution(normalized_command)
            
            # Execute the command
            process = await asyncio.create_subprocess_shell(
                normalized_command,
                stdout=asyncio.subprocess.PIPE,
                stderr=asyncio.subprocess.PIPE,
                cwd=self.workspace_root
            )
            
            try:
                stdout, stderr = await asyncio.wait_for(
                    process.communicate(),
                    timeout=timeout
                )
                
                output = stdout.decode('utf-8', errors='ignore') if stdout else ""
                error = stderr.decode('utf-8', errors='ignore') if stderr else ""
                
                # Format result for visibility in chat
                result = f"Command executed: {normalized_command}\n"
                result += f"Working directory: {self.workspace_root}\n"
                result += f"Exit code: {process.returncode}\n\n"
                
                if output:
                    result += f"STDOUT:\n{output}\n"
                if error:
                    result += f"STDERR:\n{error}\n"
                
                # If no output and no error, indicate success
                if not output and not error:
                    result += "Command completed successfully (no output).\n"
                
                self._invalidate_structure_caches()
                return [TextContent(type="text", text=result)]
            except asyncio.TimeoutError:
                process.kill()
                return [TextContent(
                    type="text",
                    text=f"Command timed out after {timeout} seconds: {normalized_command}"
                )]
        except Exception as e:
            logger.exception(f"Error executing command: {command}")
            return [TextContent(
                type="text",
                text=f"Error executing command: {str(e)}"
            )]
    
    async def _download_file(self, url: str, path: Optional[str] = None, timeout: int = 60) -> List[TextContent]:
        """Download a file from a URL and save it to the workspace"""
        if not self.file_service:
            return [TextContent(type="text", text="File service not available")]
        
        try:
            # Validate URL
            parsed_url = urlparse(url)
            if parsed_url.scheme not in ('http', 'https'):
                return [TextContent(type="text", text=f"ERROR: Invalid URL scheme. Only HTTP and HTTPS are supported. Got: {parsed_url.scheme}")]
            
            # Determine save path
            if not path:
                # Extract filename from URL
                filename = os.path.basename(parsed_url.path) or "downloaded_file"
                # Remove query parameters from filename
                if '?' in filename:
                    filename = filename.split('?')[0]
                if not filename or filename == '/':
                    filename = "downloaded_file"
                path = filename
            
            # Normalize path
            if not os.path.isabs(path):
                path = os.path.join(self.workspace_root, path)
            
            # Create directory if it doesn't exist
            directory = os.path.dirname(path)
            if directory and not os.path.exists(directory):
                os.makedirs(directory, exist_ok=True)
            
            # Download the file
            logger.info(f"Downloading file from {url} to {path}")
            async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=timeout)) as session:
                async with session.get(url) as response:
                    if response.status != 200:
                        return [TextContent(
                            type="text",
                            text=f"ERROR: Failed to download file. HTTP status: {response.status}"
                        )]
                    
                    # Check content length if available
                    content_length = response.headers.get('Content-Length')
                    if content_length:
                        size_mb = int(content_length) / (1024 * 1024)
                        if size_mb > 100:  # Limit to 100MB
                            return [TextContent(
                                type="text",
                                text=f"ERROR: File too large ({size_mb:.1f}MB). Maximum size is 100MB."
                            )]
                    
                    # Read content in chunks and write to file
                    content = await response.read()
                    
                    # Check size after download
                    size_mb = len(content) / (1024 * 1024)
                    if size_mb > 100:
                        return [TextContent(
                            type="text",
                            text=f"ERROR: Downloaded file too large ({size_mb:.1f}MB). Maximum size is 100MB."
                        )]
                    
                    # Write file
                    async with aiofiles.open(path, 'wb') as f:
                        await f.write(content)
                    
                    self._invalidate_structure_caches()
                    file_size_kb = len(content) / 1024
                    return [TextContent(
                        type="text",
                        text=f"Successfully downloaded file from {url}\nSaved to: {path}\nSize: {file_size_kb:.1f} KB"
                    )]
        except aiohttp.ClientError as e:
            return [TextContent(type="text", text=f"ERROR: Network error downloading file: {str(e)}")]
        except asyncio.TimeoutError:
            return [TextContent(type="text", text=f"ERROR: Download timed out after {timeout} seconds")]
        except Exception as e:
            logger.error(f"Error downloading file: {e}", exc_info=True)
            return [TextContent(type="text", text=f"ERROR: Failed to download file: {str(e)}")]
    
    async def _web_search(self, query: str, max_results: int, search_type: str = "text") -> List[TextContent]:
        """Perform web search using enhanced web search service"""
        try:
            from .web_search_service import WebSearchService
            
            # Use shared web search service instance if available, otherwise create new one
            if hasattr(self, '_web_search_service') and self._web_search_service:
                web_service = self._web_search_service
            else:
                web_service = WebSearchService()
                self._web_search_service = web_service
            
            # Detect if this is a price query - these need fresh, uncached results
            query_lower = query.lower()
            is_price_query = any(keyword in query_lower for keyword in [
                "price", "cost", "value", "worth", "rate", "bitcoin", "btc", "ethereum", "eth",
                "crypto", "stock", "currency", "usd", "eur", "gbp", "jpy", "exchange rate"
            ])
            
            # For price queries, disable cache and increase results to get more sources
            use_cache = not is_price_query  # Don't cache price queries
            
            # Clear any existing cache entries for price queries to ensure fresh results
            if is_price_query and hasattr(web_service, 'cache'):
                # Remove any cached entries that might match this price query
                # This ensures we always get fresh price data
                cache_keys_to_remove = []
                for cache_key in list(web_service.cache.keys()):
                    # Check if this cache key is related to price queries
                    cache_key_lower = cache_key.lower()
                    if any(term in cache_key_lower for term in [
                        "price", "bitcoin", "btc", "ethereum", "eth", "crypto", 
                        "stock", "currency", "exchange rate", "current", "live"
                    ]):
                        cache_keys_to_remove.append(cache_key)
                for key in cache_keys_to_remove:
                    web_service.cache.pop(key, None)
                # Also clear the cache file if it exists
                if hasattr(web_service, 'clear_cache'):
                    # Don't clear all cache, just price-related entries
                    pass
            
            search_max_results = max_results * 2 if is_price_query else max_results  # Get more results for prices
            
            # Optimize query for price queries to get current data
            if is_price_query:
                # Add "current" or "live" if not already present
                optimized_query = query
                if "current" not in query_lower and "live" not in query_lower and "today" not in query_lower:
                    optimized_query = f"current {query}"
                else:
                    optimized_query = query
            else:
                optimized_query = query
            
            # Perform search with enhanced features
            results, metadata = await web_service.search(
                query=optimized_query,
                max_results=min(search_max_results, 20),  # Cap at 20
                search_type=search_type,
                use_cache=use_cache,
                optimize_query=not is_price_query  # Don't optimize further if we already optimized
            )
            
            if metadata.get("error"):
                return [TextContent(type="text", text=f"Search error: {metadata['error']}")]
            
            # For price queries, prioritize results from reputable sources
            if is_price_query and results:
                # Re-sort to prioritize exchanges and financial sites
                def price_source_priority(result):
                    url = (result.get("href") or result.get("url") or "").lower()
                    title = (result.get("title") or "").lower()
                    body = (result.get("body") or result.get("description") or "").lower()
                    
                    priority = 0
                    # High priority sources
                    if any(domain in url for domain in [
                        "coinbase", "binance", "kraken", "gemini", "bitstamp", 
                        "coindesk", "cointelegraph", "bloomberg", "reuters",
                        "yahoo finance", "marketwatch", "nasdaq", "investing.com"
                    ]):
                        priority += 10
                    # Medium priority
                    elif any(domain in url for domain in [
                        "cryptocurrency", "crypto", "exchange", "trading"
                    ]):
                        priority += 5
                    # Check for price indicators in title/body
                    if any(indicator in title or indicator in body for indicator in [
                        "$", "usd", "eur", "price", "current", "live", "now"
                    ]):
                        priority += 3
                    
                    return priority
                
                # Sort by priority (highest first)
                results.sort(key=price_source_priority, reverse=True)
            
            # Format results
            formatted = web_service.format_results(results, optimized_query, include_metadata=True)
            
            # Add metadata info if cached (but price queries shouldn't be cached)
            if metadata.get("cached") and not is_price_query:
                cache_age = metadata.get("cache_age_seconds", 0)
                formatted += f"\n\n[Note: Results from cache, age: {cache_age}s]"
            elif is_price_query:
                formatted += f"\n\n[Note: Fresh search results for current price data]"
            
            return [TextContent(type="text", text=formatted)]
        except ImportError:
            return [TextContent(type="text", text="Web search not available. Install the 'ddgs' package.")]
        except Exception as e:
            return [TextContent(type="text", text=f"Error performing web search: {str(e)}")]
    
    async def _create_document(self, path: str, content: str, title: str = "", author: Optional[str] = None) -> List[TextContent]:
        """Create a Word document (.docx) with improved markdown parsing"""
        try:
            try:
                from docx import Document
                from docx.shared import Inches, Pt, RGBColor
                from docx.enum.text import WD_ALIGN_PARAGRAPH
                from docx.oxml.ns import qn
            except ImportError:
                return [TextContent(
                    type="text",
                    text="ERROR: python-docx package not installed. Install it with: pip install python-docx"
                )]
            
            # Normalize path
            if not os.path.isabs(path):
                path = os.path.join(self.workspace_root, path)
            
            # Ensure .docx extension
            if not path.lower().endswith('.docx'):
                path += '.docx'
            
            # Create directory if needed
            directory = os.path.dirname(path)
            if directory and not os.path.exists(directory):
                os.makedirs(directory, exist_ok=True)
            
            # Create document
            doc = Document()
            
            # Set document properties
            if title:
                doc.core_properties.title = title
            if author:
                doc.core_properties.author = author or "AI Agent"
            
            # Add title if provided
            if title:
                title_para = doc.add_heading(title, 0)
                title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Helper function to parse inline markdown formatting (bold, italic)
            def add_formatted_text(paragraph, text: str):
                """Add text with markdown formatting (bold **text**, italic *text*)"""
                if not text:
                    return
                
                # Process bold (**text**) and italic (*text*) together
                # First, mark bold regions, then process italic within non-bold regions
                
                # Find all bold markers (**text**)
                bold_pattern = r'\*\*([^*]+)\*\*'
                bold_matches = list(re.finditer(bold_pattern, text))
                
                # Find all italic markers (*text*) but not **text**
                # We need to avoid italic markers that are part of bold markers
                italic_pattern = r'(?<!\*)\*([^*]+?)\*(?!\*)'
                italic_matches = list(re.finditer(italic_pattern, text))
                
                # Create a list of all formatting regions
                regions = []
                for match in bold_matches:
                    regions.append({
                        'start': match.start(),
                        'end': match.end(),
                        'type': 'bold',
                        'text': match.group(1)
                    })
                for match in italic_matches:
                    # Check if this italic is inside a bold region
                    is_inside_bold = any(
                        r['start'] <= match.start() < r['end'] or 
                        r['start'] < match.end() <= r['end']
                        for r in regions if r['type'] == 'bold'
                    )
                    if not is_inside_bold:
                        regions.append({
                            'start': match.start(),
                            'end': match.end(),
                            'type': 'italic',
                            'text': match.group(1)
                        })
                
                # Sort regions by start position
                regions.sort(key=lambda x: x['start'])
                
                # Build text runs
                if regions:
                    pos = 0
                    for region in regions:
                        # Add text before this region
                        if region['start'] > pos:
                            plain_text = text[pos:region['start']]
                            if plain_text:
                                paragraph.add_run(plain_text)
                        
                        # Add formatted text
                        run = paragraph.add_run(region['text'])
                        if region['type'] == 'bold':
                            run.bold = True
                        elif region['type'] == 'italic':
                            run.italic = True
                        
                        pos = region['end']
                    
                    # Add remaining text
                    if pos < len(text):
                        plain_text = text[pos:]
                        if plain_text:
                            paragraph.add_run(plain_text)
                else:
                    # No formatting found, add text as-is (with markers removed)
                    cleaned_text = re.sub(r'\*\*([^*]+)\*\*', r'\1', text)
                    cleaned_text = re.sub(r'(?<!\*)\*([^*]+?)\*(?!\*)', r'\1', cleaned_text)
                    if cleaned_text:
                        paragraph.add_run(cleaned_text)
            
            # Parse content and add to document
            lines = content.split('\n')
            i = 0
            
            while i < len(lines):
                line = lines[i].rstrip()  # Keep leading spaces, remove trailing
                
                # Check for markdown table
                if '|' in line and line.strip().startswith('|') and line.strip().endswith('|'):
                    # Collect table rows
                    table_rows = []
                    header_row = None
                    separator_row = None
                    
                    # Check if this is a table
                    while i < len(lines) and '|' in lines[i]:
                        current_line = lines[i].strip()
                        if current_line.startswith('|') and current_line.endswith('|'):
                            # Check if it's a separator row (contains dashes and colons)
                            if re.match(r'^\|[\s\-\|:]+\|$', current_line):
                                separator_row = current_line
                            else:
                                if header_row is None:
                                    header_row = current_line
                                else:
                                    table_rows.append(current_line)
                        i += 1
                    
                    # Create table if we have a valid header
                    if header_row:
                        # Parse header
                        header_cells = [cell.strip() for cell in header_row.split('|')[1:-1]]
                        num_cols = len(header_cells)
                        
                        if num_cols > 0:
                            # Create Word table (1 row for header, will add data rows if any)
                            table = doc.add_table(rows=1, cols=num_cols)
                            table.style = 'Light Grid Accent 1'
                            
                            # Add header row
                            header_cells_row = table.rows[0].cells
                            for j, cell_text in enumerate(header_cells):
                                cell = header_cells_row[j]
                                cell.paragraphs[0].clear()
                                add_formatted_text(cell.paragraphs[0], cell_text)
                                # Make header bold
                                for paragraph in cell.paragraphs:
                                    for run in paragraph.runs:
                                        run.bold = True
                            
                            # Add data rows if any
                            for row_text in table_rows:
                                row_cells = [cell.strip() for cell in row_text.split('|')[1:-1]]
                                # Ensure we have the right number of cells
                                while len(row_cells) < num_cols:
                                    row_cells.append('')
                                row_cells = row_cells[:num_cols]
                                
                                new_row = table.add_row()
                                for j, cell_text in enumerate(row_cells):
                                    cell = new_row.cells[j]
                                    cell.paragraphs[0].clear()
                                    add_formatted_text(cell.paragraphs[0], cell_text)
                    
                    continue
                
                # Handle empty lines
                if not line.strip():
                    doc.add_paragraph()  # Empty paragraph for spacing
                    i += 1
                    continue
                
                # Handle markdown-style headings
                if line.startswith('# '):
                    heading_text = line[2:].strip()
                    para = doc.add_heading('', level=1)  # Create empty heading to preserve style
                    add_formatted_text(para, heading_text)
                elif line.startswith('## '):
                    heading_text = line[3:].strip()
                    para = doc.add_heading('', level=2)
                    add_formatted_text(para, heading_text)
                elif line.startswith('### '):
                    heading_text = line[4:].strip()
                    para = doc.add_heading('', level=3)
                    add_formatted_text(para, heading_text)
                elif line.startswith('#### '):
                    heading_text = line[5:].strip()
                    para = doc.add_heading('', level=4)
                    add_formatted_text(para, heading_text)
                elif line.startswith('##### '):
                    heading_text = line[6:].strip()
                    para = doc.add_heading('', level=5)
                    add_formatted_text(para, heading_text)
                elif line.startswith('###### '):
                    heading_text = line[7:].strip()
                    para = doc.add_heading('', level=6)
                    add_formatted_text(para, heading_text)
                # Handle bullet lists
                elif line.startswith('- ') or line.startswith('* ') or line.startswith('• '):
                    list_text = line[2:].strip()
                    para = doc.add_paragraph(style='List Bullet')
                    add_formatted_text(para, list_text)
                # Handle numbered lists
                elif re.match(r'^\d+[\.\)]\s+', line):
                    match = re.match(r'^(\d+)[\.\)]\s+(.*)$', line)
                    if match:
                        list_text = match.group(2)
                        para = doc.add_paragraph(style='List Number')
                        add_formatted_text(para, list_text)
                else:
                    # Regular paragraph
                    para = doc.add_paragraph()
                    add_formatted_text(para, line)
                
                i += 1
            
            # Save document
            doc.save(path)
            self._invalidate_structure_caches()
            
            return [TextContent(
                type="text",
                text=f"Successfully created document: {path}\nTitle: {title or 'Untitled'}\nContent length: {len(content)} characters"
            )]
        except Exception as e:
            logger.error(f"Error creating document: {e}", exc_info=True)
            return [TextContent(type="text", text=f"ERROR: Failed to create document: {str(e)}")]
    
    async def _create_slide(self, path: str, title: str, content: str = "", layout: str = "title_content") -> List[TextContent]:
        """Create a PowerPoint slide with enhanced formatting, design, backgrounds, and professional styling"""
        try:
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
                from pptx.dml.color import RGBColor
                from pptx.enum.shapes import MSO_SHAPE
            except ImportError:
                return [TextContent(
                    type="text",
                    text="ERROR: python-pptx package not installed. Install it with: pip install python-pptx"
                )]
            
            # Normalize path
            if not os.path.isabs(path):
                path = os.path.join(self.workspace_root, path)
            
            # Ensure .pptx extension
            if not path.lower().endswith('.pptx'):
                path += '.pptx'
            
            # Create directory if needed
            directory = os.path.dirname(path)
            if directory and not os.path.exists(directory):
                os.makedirs(directory, exist_ok=True)
            
            # Create or load presentation
            if os.path.exists(path):
                prs = Presentation(path)
            else:
                prs = Presentation()
                # Set slide size to widescreen (16:9) for modern look
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(5.625)
            
            # Map layout names to slide layout indices
            layout_map = {
                "title_slide": 0,
                "title_content": 1,
                "title_only": 5,
                "blank": 6
            }
            layout_idx = layout_map.get(layout, 1)  # Default to title_content
            
            # Add slide
            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)
            
            # === ENHANCED BACKGROUND DESIGN ===
            # Add gradient background rectangle covering entire slide
            background = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
            )
            background.fill.solid()
            # Professional gradient-like color (light blue-gray)
            background.fill.fore_color.rgb = RGBColor(245, 247, 250)
            background.line.fill.background()  # No border
            # Send to back
            slide.shapes._spTree.remove(background._element)
            slide.shapes._spTree.insert(2, background._element)
            
            # Add decorative accent bar at top
            accent_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
            )
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = RGBColor(31, 56, 100)  # Professional blue
            accent_bar.line.fill.background()
            slide.shapes._spTree.remove(accent_bar._element)
            slide.shapes._spTree.insert(2, accent_bar._element)
            
            # Add subtle decorative element (bottom right corner)
            decor_shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                prs.slide_width - Inches(2.5),
                prs.slide_height - Inches(1.2),
                Inches(2.5),
                Inches(1.2)
            )
            decor_shape.fill.solid()
            decor_shape.fill.fore_color.rgb = RGBColor(230, 235, 240)  # Very light gray
            decor_shape.line.fill.background()
            decor_shape.rotation = -5  # Slight rotation for visual interest
            slide.shapes._spTree.remove(decor_shape._element)
            slide.shapes._spTree.insert(2, decor_shape._element)
            
            # === ENHANCED TITLE STYLING ===
            if slide.shapes.title:
                title_shape = slide.shapes.title
                # Clear any default placeholder text first
                if hasattr(title_shape, 'text_frame'):
                    title_shape.text_frame.clear()
                title_shape.text = title
                # Enhanced title styling
                if hasattr(title_shape, 'text_frame') and len(title_shape.text_frame.paragraphs) > 0:
                    title_paragraph = title_shape.text_frame.paragraphs[0]
                    title_paragraph.font.size = Pt(48)  # Larger, more prominent
                    title_paragraph.font.bold = True
                    title_paragraph.font.color.rgb = RGBColor(31, 56, 100)  # Professional dark blue
                    title_paragraph.alignment = PP_ALIGN.LEFT
                    title_paragraph.space_after = Pt(12)
                    
                    # Adjust title position for better spacing
                    title_shape.left = Inches(0.5)
                    title_shape.top = Inches(0.3)
                    title_shape.width = Inches(9)
                    title_shape.height = Inches(1.2)
                    
                    # Add subtle shadow effect by creating a duplicate slightly offset
                    # (Note: python-pptx doesn't support shadows directly, but we can enhance with background)
                    title_shape.text_frame.vertical_anchor = MSO_ANCHOR.TOP
            
            # Helper function to add formatted content with enhanced styling
            def add_formatted_content(text_frame, content_text):
                """Add content with proper formatting, enhanced styling, and visual hierarchy"""
                text_frame.clear()
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.2)
                text_frame.margin_right = Inches(0.2)
                text_frame.margin_top = Inches(0.15)
                text_frame.margin_bottom = Inches(0.15)
                text_frame.vertical_anchor = MSO_ANCHOR.TOP
                
                lines = content_text.split('\n')
                for i, line in enumerate(lines):
                    line = line.strip()
                    if not line:
                        if i < len(lines) - 1:  # Don't add empty paragraph at the end
                            p = text_frame.add_paragraph()
                            p.space_after = Pt(8)
                        continue
                    
                    # Check for bullet points with enhanced styling
                    if line.startswith('- ') or line.startswith('* ') or line.startswith('• '):
                        bullet_text = line[2:].strip()
                        p = text_frame.add_paragraph()
                        p.text = bullet_text
                        p.level = 0
                        p.font.size = Pt(20)  # Slightly larger for readability
                        p.font.color.rgb = RGBColor(45, 45, 45)  # Darker for better contrast
                        p.space_after = Pt(10)
                        p.space_before = Pt(4)
                        p.line_spacing = 1.2
                        
                        # Enhanced formatting for key-value pairs
                        if ':' in bullet_text:
                            parts = bullet_text.split(':', 1)
                            if len(parts) == 2:
                                p.clear()
                                # Key part (bold, colored)
                                run1 = p.add_run()
                                run1.text = parts[0] + ':'
                                run1.font.bold = True
                                run1.font.size = Pt(20)
                                run1.font.color.rgb = RGBColor(31, 56, 100)  # Professional blue
                                # Value part
                                run2 = p.add_run()
                                run2.text = ' ' + parts[1]
                                run2.font.bold = False
                                run2.font.size = Pt(20)
                                run2.font.color.rgb = RGBColor(45, 45, 45)
                    elif line.startswith('  - ') or line.startswith('  * ') or line.startswith('    - '):
                        # Nested bullet with enhanced styling
                        bullet_text = line.lstrip(' -*•').strip()
                        p = text_frame.add_paragraph()
                        p.text = bullet_text
                        p.level = 1
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(60, 60, 60)
                        p.space_after = Pt(8)
                        p.space_before = Pt(2)
                        p.line_spacing = 1.15
                    elif line.startswith('## '):
                        # Subheading with enhanced styling
                        heading_text = line[3:].strip()
                        p = text_frame.add_paragraph()
                        p.text = heading_text
                        p.font.size = Pt(26)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(31, 56, 100)
                        p.space_after = Pt(14)
                        p.space_before = Pt(16)
                        p.line_spacing = 1.3
                    elif line.startswith('# '):
                        # Main heading with enhanced styling
                        heading_text = line[2:].strip()
                        p = text_frame.add_paragraph()
                        p.text = heading_text
                        p.font.size = Pt(32)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(31, 56, 100)
                        p.space_after = Pt(18)
                        p.space_before = Pt(20)
                        p.line_spacing = 1.3
                    else:
                        # Regular paragraph with enhanced styling
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.font.size = Pt(20)
                        p.font.color.rgb = RGBColor(45, 45, 45)
                        p.space_after = Pt(12)
                        p.line_spacing = 1.25
                        # Enhanced bold text formatting (**text**)
                        if '**' in line:
                            parts = line.split('**')
                            p.clear()
                            for j, part in enumerate(parts):
                                if j % 2 == 0:
                                    run = p.add_run()
                                    run.text = part
                                    run.font.bold = False
                                    run.font.size = Pt(20)
                                    run.font.color.rgb = RGBColor(45, 45, 45)
                                else:
                                    run = p.add_run()
                                    run.text = part
                                    run.font.bold = True
                                    run.font.size = Pt(20)
                                    run.font.color.rgb = RGBColor(31, 56, 100)  # Professional blue for emphasis
            
            # Clear any default placeholder text first
            for shape in slide.placeholders:
                if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                    try:
                        # Clear any default placeholder text
                        if shape.text_frame.text.strip():
                            shape.text_frame.clear()
                    except Exception:
                        pass  # Ignore errors when clearing placeholders
            
            # === ENHANCED CONTENT POSITIONING ===
            if content:
                content_placeholder = None
                # Try to find content placeholder (usually idx 1 for content)
                for shape in slide.placeholders:
                    try:
                        if (hasattr(shape, 'placeholder_format') and 
                            hasattr(shape.placeholder_format, 'idx') and
                            shape.placeholder_format.idx == 1 and 
                            shape != slide.shapes.title):
                            content_placeholder = shape
                            break
                    except Exception:
                        continue
                
                if content_placeholder and hasattr(content_placeholder, 'text_frame'):
                    # Clear any existing text first
                    content_placeholder.text_frame.clear()
                    # Enhanced positioning for placeholder
                    content_placeholder.left = Inches(0.6)
                    content_placeholder.top = Inches(1.8)
                    content_placeholder.width = Inches(8.8)
                    content_placeholder.height = Inches(3.5)
                    add_formatted_content(content_placeholder.text_frame, content)
                else:
                    # Add enhanced text box with professional positioning
                    left = Inches(0.6)
                    top = Inches(1.8)
                    width = Inches(8.8)
                    height = Inches(3.5)
                    
                    # Add subtle background box for content area
                    content_bg = slide.shapes.add_shape(
                        MSO_SHAPE.ROUNDED_RECTANGLE,
                        left - Inches(0.1),
                        top - Inches(0.1),
                        width + Inches(0.2),
                        height + Inches(0.2)
                    )
                    content_bg.fill.solid()
                    content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
                    content_bg.line.color.rgb = RGBColor(220, 225, 230)  # Light border
                    content_bg.line.width = Pt(1)
                    # Send background to back
                    slide.shapes._spTree.remove(content_bg._element)
                    slide.shapes._spTree.insert(2, content_bg._element)
                    
                    # Add text box on top of background
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    text_frame.clear()
                    add_formatted_content(text_frame, content)
            
            # Save presentation
            prs.save(path)
            self._invalidate_structure_caches()
            
            # Return a clear status message that won't be interpreted as file content
            # Use a clear prefix to indicate this is a tool result, not file content
            return [TextContent(
                type="text",
                text=f"[TOOL RESULT - DO NOT SAVE AS FILE]\n\n✅ PowerPoint slide created successfully with enhanced design!\n\n📄 File saved to: {path}\n📋 Slide title: {title}\n🎨 Layout: {layout}\n\n✨ Design Enhancements Applied:\n  • Professional gradient background\n  • Decorative accent bar\n  • Enhanced typography and spacing\n  • Improved content positioning\n  • Visual hierarchy optimization\n  • Professional color scheme\n\n⚠️ IMPORTANT: This is a status message from the create_slide tool. The .pptx file has already been created. Do NOT create any text files or save this message as file content."
            )]
        except Exception as e:
            logger.error(f"Error creating slide: {e}", exc_info=True)
            return [TextContent(type="text", text=f"ERROR: Failed to create slide: {str(e)}")]
    
    async def _create_presentation(self, path: str, title: str, slides: List[Dict[str, Any]], author: Optional[str] = None) -> List[TextContent]:
        """Create a full PowerPoint presentation with enhanced formatting and design"""
        try:
            try:
                from pptx import Presentation
                from pptx.util import Inches, Pt
                from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
                from pptx.dml.color import RGBColor
            except ImportError:
                return [TextContent(
                    type="text",
                    text="ERROR: python-pptx package not installed. Install it with: pip install python-pptx"
                )]
            
            # Normalize path
            if not os.path.isabs(path):
                path = os.path.join(self.workspace_root, path)
            
            # Ensure .pptx extension
            if not path.lower().endswith('.pptx'):
                path += '.pptx'
            
            # Create directory if needed
            directory = os.path.dirname(path)
            if directory and not os.path.exists(directory):
                os.makedirs(directory, exist_ok=True)
            
            # Create presentation
            prs = Presentation()
            
            # Set presentation properties
            if title:
                prs.core_properties.title = title
            if author:
                prs.core_properties.author = author or "AI Agent"
            
            # Helper function to add formatted content (same as in _create_slide)
            def add_formatted_content(text_frame, content_text):
                """Add content with proper formatting (bullet points, paragraphs, etc.)"""
                text_frame.clear()
                text_frame.word_wrap = True
                text_frame.margin_bottom = Inches(0.1)
                
                lines = content_text.split('\n')
                for i, line in enumerate(lines):
                    line = line.strip()
                    if not line:
                        if i < len(lines) - 1:  # Don't add empty paragraph at the end
                            p = text_frame.add_paragraph()
                            p.space_after = Pt(6)
                        continue
                    
                    # Check for bullet points
                    if line.startswith('- ') or line.startswith('* ') or line.startswith('• '):
                        bullet_text = line[2:].strip()
                        p = text_frame.add_paragraph()
                        p.text = bullet_text
                        p.level = 0
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(51, 51, 51)  # Dark gray
                        p.space_after = Pt(8)
                        # Make first word bold if it looks like a heading
                        if ':' in bullet_text:
                            parts = bullet_text.split(':', 1)
                            if len(parts) == 2:
                                run = p.runs[0]
                                run.text = parts[0] + ':'
                                run.font.bold = True
                                run2 = p.add_run()
                                run2.text = parts[1]
                                run2.font.bold = False
                    elif line.startswith('  - ') or line.startswith('  * ') or line.startswith('    - '):
                        # Nested bullet
                        bullet_text = line.lstrip(' -*•').strip()
                        p = text_frame.add_paragraph()
                        p.text = bullet_text
                        p.level = 1
                        p.font.size = Pt(16)
                        p.font.color.rgb = RGBColor(68, 68, 68)
                        p.space_after = Pt(6)
                    elif line.startswith('## '):
                        # Subheading
                        heading_text = line[3:].strip()
                        p = text_frame.add_paragraph()
                        p.text = heading_text
                        p.font.size = Pt(24)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(31, 56, 100)
                        p.space_after = Pt(12)
                        p.space_before = Pt(12)
                    elif line.startswith('# '):
                        # Main heading
                        heading_text = line[2:].strip()
                        p = text_frame.add_paragraph()
                        p.text = heading_text
                        p.font.size = Pt(28)
                        p.font.bold = True
                        p.font.color.rgb = RGBColor(31, 56, 100)
                        p.space_after = Pt(16)
                        p.space_before = Pt(16)
                    else:
                        # Regular paragraph
                        p = text_frame.add_paragraph()
                        p.text = line
                        p.font.size = Pt(18)
                        p.font.color.rgb = RGBColor(51, 51, 51)
                        p.space_after = Pt(10)
                        # Check for bold text (**text**)
                        if '**' in line:
                            parts = line.split('**')
                            p.clear()
                            for j, part in enumerate(parts):
                                if j % 2 == 0:
                                    run = p.add_run()
                                    run.text = part
                                    run.font.bold = False
                                else:
                                    run = p.add_run()
                                    run.text = part
                                    run.font.bold = True
                                    run.font.size = Pt(18)
                                    run.font.color.rgb = RGBColor(31, 56, 100)
            
            # Add title slide with enhanced styling
            title_slide_layout = prs.slide_layouts[0]
            title_slide = prs.slides.add_slide(title_slide_layout)
            if title_slide.shapes.title:
                title_shape = title_slide.shapes.title
                # Clear any default placeholder text first
                if hasattr(title_shape, 'text_frame'):
                    title_shape.text_frame.clear()
                title_shape.text = title
                # Style the title
                if hasattr(title_shape, 'text_frame') and len(title_shape.text_frame.paragraphs) > 0:
                    title_paragraph = title_shape.text_frame.paragraphs[0]
                    title_paragraph.font.size = Pt(54)
                    title_paragraph.font.bold = True
                    title_paragraph.font.color.rgb = RGBColor(31, 56, 100)
                    title_paragraph.alignment = PP_ALIGN.CENTER
            
            # Add subtitle if author is provided
            if author and len(title_slide.shapes) > 1:
                subtitle_shape = title_slide.shapes[1]
                if hasattr(subtitle_shape, 'text_frame'):
                    # Clear any default placeholder text
                    subtitle_shape.text_frame.clear()
                    subtitle_shape.text = f"Created by {author}"
                    if len(subtitle_shape.text_frame.paragraphs) > 0:
                        subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
                        subtitle_paragraph.font.size = Pt(24)
                        subtitle_paragraph.font.color.rgb = RGBColor(100, 100, 100)
                        subtitle_paragraph.alignment = PP_ALIGN.CENTER
            
            # Add content slides
            layout_map = {
                "title_slide": 0,
                "title_content": 1,
                "title_only": 5,
                "blank": 6
            }
            
            for slide_data in slides:
                slide_title = slide_data.get("title", "")
                slide_content = slide_data.get("content", "")
                slide_layout_name = slide_data.get("layout", "title_content")
                layout_idx = layout_map.get(slide_layout_name, 1)
                
                slide_layout = prs.slide_layouts[layout_idx]
                slide = prs.slides.add_slide(slide_layout)
                
                # Set title with enhanced styling
                if slide.shapes.title:
                    title_shape = slide.shapes.title
                    # Clear any default placeholder text first
                    if hasattr(title_shape, 'text_frame'):
                        title_shape.text_frame.clear()
                    title_shape.text = slide_title
                    # Style the title
                    if hasattr(title_shape, 'text_frame') and len(title_shape.text_frame.paragraphs) > 0:
                        title_paragraph = title_shape.text_frame.paragraphs[0]
                        title_paragraph.font.size = Pt(44)
                        title_paragraph.font.bold = True
                        title_paragraph.font.color.rgb = RGBColor(31, 56, 100)
                        title_paragraph.alignment = PP_ALIGN.LEFT
                
                # Clear any default placeholder text first
                for shape in slide.placeholders:
                    if hasattr(shape, 'text_frame') and shape != slide.shapes.title:
                        try:
                            # Clear any default placeholder text
                            if shape.text_frame.text.strip():
                                shape.text_frame.clear()
                        except Exception:
                            pass  # Ignore errors when clearing placeholders
                
                # Set content with formatting
                if slide_content:
                    content_placeholder = None
                    # Try to find content placeholder (usually idx 1 for content)
                    for shape in slide.placeholders:
                        try:
                            if (hasattr(shape, 'placeholder_format') and 
                                hasattr(shape.placeholder_format, 'idx') and
                                shape.placeholder_format.idx == 1 and 
                                shape != slide.shapes.title):
                                content_placeholder = shape
                                break
                        except Exception:
                            continue
                    
                    if content_placeholder and hasattr(content_placeholder, 'text_frame'):
                        # Clear any existing text first
                        content_placeholder.text_frame.clear()
                        add_formatted_content(content_placeholder.text_frame, slide_content)
                    else:
                        # Add text box with better positioning
                        left = Inches(0.7)
                        top = Inches(2.2)
                        width = Inches(8.6)
                        height = Inches(4.3)
                        text_box = slide.shapes.add_textbox(left, top, width, height)
                        text_frame = text_box.text_frame
                        text_frame.clear()  # Clear any default text
                        add_formatted_content(text_frame, slide_content)
            
            # Save presentation
            prs.save(path)
            self._invalidate_structure_caches()
            
            # Return a clear status message that won't be interpreted as file content
            # Use a clear prefix to indicate this is a tool result, not file content
            return [TextContent(
                type="text",
                text=f"[TOOL RESULT - DO NOT SAVE AS FILE]\n\n✅ PowerPoint presentation created successfully!\n\n📄 File saved to: {path}\n📋 Presentation title: {title}\n📊 Total slides: {len(slides) + 1} (including title slide)\n✨ Enhanced formatting and design applied\n\n⚠️ IMPORTANT: This is a status message from the create_presentation tool. The .pptx file has already been created. Do NOT create any text files or save this message as file content."
            )]
        except Exception as e:
            logger.error(f"Error creating presentation: {e}", exc_info=True)
            return [TextContent(type="text", text=f"ERROR: Failed to create presentation: {str(e)}")]
    
    async def _get_user_location(self, ip_address: Optional[str] = None) -> List[TextContent]:
        """Get user location based on IP address"""
        if not self._location_service:
            return [TextContent(
                type="text",
                text="Location service not available. Please ensure location_service is initialized."
            )]
        
        try:
            location = await self._location_service.get_user_location(ip_address)
            
            if "error" in location:
                return [TextContent(
                    type="text",
                    text=f"Error getting location: {location.get('error', 'Unknown error')}"
                )]
            
            # Format location information
            result = "User Location Information:\n\n"
            result += f"City: {location.get('city', 'Unknown')}\n"
            result += f"Region: {location.get('region', 'Unknown')}\n"
            result += f"Country: {location.get('country', 'Unknown')} ({location.get('country_code', '')})\n"
            result += f"Coordinates: {location.get('latitude', 0):.4f}, {location.get('longitude', 0):.4f}\n"
            result += f"Timezone: {location.get('timezone', 'Unknown')}\n"
            result += f"IP Address: {location.get('ip', 'Unknown')}\n"
            
            return [TextContent(type="text", text=result)]
        except Exception as e:
            logger.error(f"Error getting user location: {e}", exc_info=True)
            return [TextContent(type="text", text=f"Error getting location: {str(e)}")]
    
    async def _get_weather(
        self,
        city: Optional[str],
        latitude: Optional[float],
        longitude: Optional[float],
        units: str
    ) -> List[TextContent]:
        """Get weather for a location"""
        if not self._location_service:
            return [TextContent(
                type="text",
                text="Location service not available. Please ensure location_service is initialized."
            )]
        
        try:
            weather = await self._location_service.get_weather(city, latitude, longitude, units)
            
            if "error" in weather:
                return [TextContent(
                    type="text",
                    text=f"Error getting weather: {weather.get('error', 'Unknown error')}"
                )]
            
            # Format weather information
            temp_unit = "°C" if units == "metric" else "°F"
            speed_unit = "km/h" if units == "metric" else "mph"
            
            result = f"Weather Information for {weather.get('location', 'Unknown')}:\n\n"
            result += f"Temperature: {weather.get('temperature', 0):.1f}{temp_unit}\n"
            result += f"Feels Like: {weather.get('feels_like', 0):.1f}{temp_unit}\n"
            result += f"Conditions: {weather.get('description', 'Unknown').title()}\n"
            result += f"Humidity: {weather.get('humidity', 0)}%\n"
            result += f"Wind Speed: {weather.get('wind_speed', 0):.1f} {speed_unit}\n"
            result += f"Pressure: {weather.get('pressure', 0)} hPa\n"
            if weather.get('visibility'):
                result += f"Visibility: {weather.get('visibility', 0) / 1000:.1f} km\n"
            result += f"\nData Source: {weather.get('source', 'Unknown')}\n"
            
            return [TextContent(type="text", text=result)]
        except Exception as e:
            logger.error(f"Error getting weather: {e}", exc_info=True)
            return [TextContent(type="text", text=f"Error getting weather: {str(e)}")]
    
    async def _get_news(
        self,
        city: Optional[str],
        country: Optional[str],
        query: Optional[str],
        max_results: int
    ) -> List[TextContent]:
        """Get news for a location"""
        if not self._location_service:
            return [TextContent(
                type="text",
                text="Location service not available. Please ensure location_service is initialized."
            )]
        
        try:
            news = await self._location_service.get_news(city, country, query, max_results)
            
            if "error" in news:
                return [TextContent(
                    type="text",
                    text=f"Error getting news: {news.get('error', 'Unknown error')}"
                )]
            
            articles = news.get("articles", [])
            if not articles:
                return [TextContent(
                    type="text",
                    text=f"No news articles found for the specified location/query."
                )]
            
            # Format news information
            result = f"News Articles ({news.get('total_results', len(articles))} results):\n\n"
            result += f"Query: {news.get('query', 'General news')}\n"
            result += f"Source: {news.get('source', 'Unknown')}\n\n"
            result += "=" * 80 + "\n\n"
            
            for i, article in enumerate(articles[:max_results], 1):
                result += f"{i}. {article.get('title', 'No title')}\n"
                if article.get('description'):
                    result += f"   {article.get('description', '')}\n"
                if article.get('source'):
                    result += f"   Source: {article.get('source', '')}\n"
                if article.get('published_at'):
                    result += f"   Published: {article.get('published_at', '')}\n"
                if article.get('url'):
                    result += f"   URL: {article.get('url', '')}\n"
                result += "\n"
            
            return [TextContent(type="text", text=result)]
        except Exception as e:
            logger.error(f"Error getting news: {e}", exc_info=True)
            return [TextContent(type="text", text=f"Error getting news: {str(e)}")]
    
    async def _save_memory(self, content: str) -> List[TextContent]:
        """Save a memory"""
        if not self._memory_service:
            return [TextContent(
                type="text",
                text="Memory service not available. Please ensure memory_service is initialized."
            )]
        
        if not content or not content.strip():
            return [TextContent(
                type="text",
                text="ERROR: Memory content cannot be empty"
            )]
        
        try:
            memory = self._memory_service.add_memory(content.strip())
            return [TextContent(
                type="text",
                text=f"✅ Memory saved successfully!\n\nSaved: {memory['content']}\n\nThis information will now be remembered for future conversations."
            )]
        except Exception as e:
            logger.error(f"Error saving memory: {e}", exc_info=True)
            return [TextContent(type="text", text=f"Error saving memory: {str(e)}")]
    
    async def _fetch_cryptocraft_prediction(self, coin_id: str, asset_name: str, days_ahead: int) -> Optional[Dict[str, Any]]:
        """Fetch price prediction from CryptoCraft.com"""
        try:
            # Map coin IDs to CryptoCraft URL format
            cryptocraft_symbols = {
                'bitcoin': 'BTC',
                'ethereum': 'ETH',
                'cardano': 'ADA',
                'solana': 'SOL',
                'polkadot': 'DOT',
                'chainlink': 'LINK',
                'avalanche': 'AVAX',
                'polygon': 'MATIC',
                'dogecoin': 'DOGE',
                'litecoin': 'LTC',
                'ripple': 'XRP',
            }
            
            symbol = cryptocraft_symbols.get(coin_id.lower())
            if not symbol:
                return None
            
            # Try to fetch from CryptoCraft.com
            # Note: This attempts to fetch prediction data from the website
            url = f"https://www.cryptocraft.com/{symbol.lower()}"
            
            async with aiohttp.ClientSession() as session:
                headers = {
                    'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
                }
                async with session.get(url, headers=headers, timeout=aiohttp.ClientTimeout(total=10)) as response:
                    if response.status == 200:
                        html_content = await response.text()
                        
                        # Try to extract prediction data from HTML
                        # Look for common patterns: price predictions, forecast values, etc.
                        import re
                        
                        # Try to find price predictions in various formats
                        # Pattern 1: Look for numbers that could be price predictions
                        # Pattern 2: Look for percentage changes
                        # Pattern 3: Look for forecast/prediction keywords near numbers
                        
                        # Extract potential price values (numbers with $ or without)
                        price_patterns = [
                            r'\$[\d,]+\.?\d*',  # $50,000 or $50,000.50
                            r'[\d,]+\.?\d*\s*(?:USD|usd|\$)',  # 50000 USD
                        ]
                        
                        # Look for prediction-related text
                        prediction_keywords = ['predict', 'forecast', 'target', 'expect', 'project']
                        prediction_sections = []
                        
                        # Split HTML into sections and look for prediction-related content
                        for keyword in prediction_keywords:
                            pattern = rf'{keyword}[^<]*?[\d,]+\.?\d*'
                            matches = re.findall(pattern, html_content, re.IGNORECASE)
                            prediction_sections.extend(matches)
                        
                        # Try to extract meaningful price predictions
                        # This is a simplified approach - in production, you'd want more sophisticated parsing
                        extracted_prices = []
                        for section in prediction_sections[:10]:  # Limit to first 10 matches
                            # Extract numbers from the section
                            numbers = re.findall(r'[\d,]+\.?\d*', section)
                            for num_str in numbers:
                                try:
                                    price_val = float(num_str.replace(',', ''))
                                    # Filter reasonable price ranges based on asset
                                    if coin_id.lower() == 'bitcoin' and 10000 < price_val < 200000:
                                        extracted_prices.append(price_val)
                                    elif coin_id.lower() == 'ethereum' and 500 < price_val < 20000:
                                        extracted_prices.append(price_val)
                                    elif 0.01 < price_val < 100000:  # General range for other cryptos
                                        extracted_prices.append(price_val)
                                except ValueError:
                                    continue
                        
                        if extracted_prices:
                            # Use median or average of extracted prices
                            import statistics
                            predicted_price = statistics.median(extracted_prices) if len(extracted_prices) > 1 else extracted_prices[0]
                            
                            # Calculate confidence based on number of matches
                            confidence = min(85, 50 + (len(extracted_prices) * 5))
                            
                            return {
                                'predicted_price': predicted_price,
                                'confidence': confidence,
                                'source': 'CryptoCraft.com',
                                'method': 'web_scraping'
                            }
            
            return None
        except Exception as e:
            logger.debug(f"Error fetching CryptoCraft prediction: {e}")
            return None
    
    async def _predict_price(
        self,
        asset: str,
        asset_type: Optional[str],
        days_ahead: int,
        include_analysis: bool
    ) -> List[TextContent]:
        """Predict future price movements using technical analysis and statistical forecasting"""
        try:
            # Import market data functions
            from backend.api.market_data import fetch_crypto_price, fetch_forex_rate
            
            asset_lower = asset.lower().strip()
            
            # Determine asset type if not provided
            if not asset_type:
                forex_pairs = ['eur/usd', 'gbp/usd', 'usd/jpy', 'usd/chf', 'aud/usd', 'usd/cad']
                if any(pair in asset_lower for pair in forex_pairs):
                    asset_type = 'forex'
                else:
                    asset_type = 'crypto'
            
            # Fetch historical data
            historical_data = None
            current_price = None
            asset_name = None
            
            if asset_type == 'crypto':
                crypto_id_map = {
                    'bitcoin': 'bitcoin', 'btc': 'bitcoin',
                    'ethereum': 'ethereum', 'eth': 'ethereum',
                    'cardano': 'cardano', 'ada': 'cardano',
                    'solana': 'solana', 'sol': 'solana',
                    'polkadot': 'polkadot', 'dot': 'polkadot',
                    'chainlink': 'chainlink', 'link': 'chainlink',
                    'avalanche': 'avalanche', 'avax': 'avalanche',
                    'polygon': 'polygon', 'matic': 'polygon',
                    'dogecoin': 'dogecoin', 'doge': 'dogecoin',
                    'litecoin': 'litecoin', 'ltc': 'litecoin',
                    'ripple': 'ripple', 'xrp': 'ripple',
                }
                coin_id = crypto_id_map.get(asset_lower, asset_lower)
                price_data = await fetch_crypto_price(coin_id, days=60)  # Get more data for better prediction
                historical_data = price_data.get('historicalData', [])
                current_price = price_data.get('currentPrice', 0)
                asset_name = price_data.get('assetName', asset.capitalize())
            elif asset_type == 'forex':
                forex_map = {
                    'eur/usd': {'base': 'EUR', 'target': 'USD'},
                    'gbp/usd': {'base': 'GBP', 'target': 'USD'},
                    'usd/jpy': {'base': 'USD', 'target': 'JPY'},
                    'usd/chf': {'base': 'USD', 'target': 'CHF'},
                    'aud/usd': {'base': 'AUD', 'target': 'USD'},
                    'usd/cad': {'base': 'USD', 'target': 'CAD'},
                }
                pair_info = forex_map.get(asset_lower)
                if pair_info:
                    price_data = await fetch_forex_rate(pair_info['base'], pair_info['target'], days=60)
                    historical_data = price_data.get('historicalData', [])
                    current_price = price_data.get('currentPrice', 0)
                    asset_name = price_data.get('assetName', f"{pair_info['base']}/{pair_info['target']}")
                else:
                    return [TextContent(type="text", text=f"Unsupported forex pair: {asset}")]
            
            if not historical_data or len(historical_data) < 7:
                return [TextContent(type="text", text=f"Insufficient historical data for {asset_name}. Need at least 7 days of data.")]
            
            if current_price == 0:
                return [TextContent(type="text", text=f"Could not fetch current price for {asset_name}")]
            
            # Extract price series
            prices = [float(item.get('price', 0)) for item in historical_data if item.get('price', 0) > 0]
            if len(prices) < 7:
                return [TextContent(type="text", text=f"Insufficient price data for analysis")]
            
            # Calculate technical indicators
            import statistics
            import math
            
            # Moving averages
            sma_7 = statistics.mean(prices[-7:]) if len(prices) >= 7 else current_price
            sma_14 = statistics.mean(prices[-14:]) if len(prices) >= 14 else current_price
            sma_30 = statistics.mean(prices[-30:]) if len(prices) >= 30 else current_price
            
            # Calculate returns and volatility
            returns = []
            for i in range(1, len(prices)):
                if prices[i-1] > 0:
                    ret = (prices[i] - prices[i-1]) / prices[i-1]
                    returns.append(ret)
            
            if not returns:
                return [TextContent(type="text", text=f"Could not calculate returns for {asset_name}")]
            
            avg_return = statistics.mean(returns)
            volatility = statistics.stdev(returns) if len(returns) > 1 else abs(avg_return)
            annualized_volatility = volatility * math.sqrt(365)  # Annualized
            
            # Calculate RSI (Relative Strength Index)
            gains = [r for r in returns if r > 0]
            losses = [-r for r in returns if r < 0]
            avg_gain = statistics.mean(gains) if gains else 0.001
            avg_loss = statistics.mean(losses) if losses else 0.001
            rs = avg_gain / avg_loss if avg_loss > 0 else 100
            rsi = 100 - (100 / (1 + rs))
            
            # Trend analysis
            recent_trend = (prices[-1] - prices[-7]) / prices[-7] if len(prices) >= 7 else 0
            short_term_trend = (prices[-1] - prices[-3]) / prices[-3] if len(prices) >= 3 else 0
            
            # Momentum
            momentum = (prices[-1] - prices[-10]) / prices[-10] if len(prices) >= 10 else 0
            
            # Price prediction using multiple methods
            predictions = []
            
            # Method 1: Trend extrapolation
            trend_prediction = current_price * (1 + recent_trend * days_ahead / 7)
            predictions.append(('Trend Extrapolation', trend_prediction))
            
            # Method 2: Moving average crossover
            if sma_7 > sma_14 > sma_30:
                # Bullish trend
                ma_prediction = current_price * (1 + abs(avg_return) * days_ahead)
            elif sma_7 < sma_14 < sma_30:
                # Bearish trend
                ma_prediction = current_price * (1 - abs(avg_return) * days_ahead)
            else:
                # Sideways
                ma_prediction = current_price * (1 + avg_return * days_ahead)
            predictions.append(('Moving Average', ma_prediction))
            
            # Method 3: Mean reversion (if RSI indicates overbought/oversold)
            if rsi > 70:
                # Overbought - expect downward correction
                reversion_prediction = current_price * (1 - volatility * days_ahead * 0.5)
            elif rsi < 30:
                # Oversold - expect upward correction
                reversion_prediction = current_price * (1 + volatility * days_ahead * 0.5)
            else:
                # Neutral
                reversion_prediction = current_price * (1 + avg_return * days_ahead)
            predictions.append(('Mean Reversion', reversion_prediction))
            
            # Method 4: Volatility-adjusted forecast
            volatility_prediction = current_price * (1 + avg_return * days_ahead)
            predictions.append(('Volatility Model', volatility_prediction))
            
            # Method 5: CryptoCraft.com prediction (for crypto only)
            cryptocraft_prediction = None
            cryptocraft_confidence = None
            if asset_type == 'crypto':
                try:
                    cryptocraft_data = await self._fetch_cryptocraft_prediction(coin_id, asset_name, days_ahead)
                    if cryptocraft_data and cryptocraft_data.get('predicted_price'):
                        cryptocraft_prediction = cryptocraft_data['predicted_price']
                        cryptocraft_confidence = cryptocraft_data.get('confidence', 0)
                        predictions.append(('CryptoCraft.com', cryptocraft_prediction))
                except Exception as e:
                    logger.debug(f"Could not fetch CryptoCraft prediction: {e}")
                    # Continue without CryptoCraft data if unavailable
            
            # Weighted average of predictions (more weight to recent trends and CryptoCraft if available)
            if cryptocraft_prediction:
                # Include CryptoCraft with higher weight if available
                weights = [0.25, 0.2, 0.2, 0.15, 0.2]  # CryptoCraft gets 20% weight
            else:
                weights = [0.3, 0.25, 0.25, 0.2]  # Original weights
            weighted_prediction = sum(pred[1] * weight for pred, weight in zip(predictions, weights))
            
            # Calculate confidence intervals
            prediction_std = volatility * current_price * math.sqrt(days_ahead)
            upper_bound = weighted_prediction + (1.96 * prediction_std)  # 95% confidence
            lower_bound = weighted_prediction - (1.96 * prediction_std)
            
            # Confidence score based on data quality and volatility
            data_quality_score = min(100, (len(prices) / 60) * 100)
            volatility_score = max(0, 100 - (annualized_volatility * 100))
            confidence_score = (data_quality_score + volatility_score) / 2
            
            # Determine trend direction
            if weighted_prediction > current_price * 1.02:
                trend_direction = "BULLISH 📈"
                trend_strength = "Strong" if abs(recent_trend) > 0.05 else "Moderate"
            elif weighted_prediction < current_price * 0.98:
                trend_direction = "BEARISH 📉"
                trend_strength = "Strong" if abs(recent_trend) > 0.05 else "Moderate"
            else:
                trend_direction = "NEUTRAL ➡️"
                trend_strength = "Sideways"
            
            # Helper function to format prices based on asset type
            def format_price(price):
                if asset_type == 'forex':
                    return f"${price:,.4f}"
                else:
                    return f"${price:,.2f}"
            
            # Format output
            result = f"📊 Price Prediction for {asset_name}\n"
            result += "=" * 60 + "\n\n"
            
            result += f"💰 Current Price: {format_price(current_price)}\n"
            result += f"📅 Prediction Period: {days_ahead} day{'s' if days_ahead != 1 else ''} ahead\n\n"
            
            result += f"🎯 Predicted Price: {format_price(weighted_prediction)}\n"
            result += f"   Expected Change: {((weighted_prediction - current_price) / current_price * 100):+.2f}%\n"
            result += f"   Trend: {trend_direction} ({trend_strength})\n\n"
            
            result += f"📈 Confidence Interval (95%):\n"
            result += f"   Upper Bound: {format_price(upper_bound)}\n"
            result += f"   Lower Bound: {format_price(lower_bound)}\n"
            result += f"   Confidence Score: {confidence_score:.1f}%\n\n"
            
            if include_analysis:
                result += "🔍 Technical Analysis:\n"
                result += "-" * 60 + "\n"
                result += f"RSI (14-day): {rsi:.2f} "
                if rsi > 70:
                    result += "(Overbought - Potential Downward Correction)\n"
                elif rsi < 30:
                    result += "(Oversold - Potential Upward Correction)\n"
                else:
                    result += "(Neutral)\n"
                
                result += f"Moving Averages:\n"
                result += f"  • 7-day SMA: {format_price(sma_7)}\n"
                result += f"  • 14-day SMA: {format_price(sma_14)}\n"
                result += f"  • 30-day SMA: {format_price(sma_30)}\n"
                
                result += f"\nVolatility Analysis:\n"
                result += f"  • Daily Volatility: {volatility*100:.2f}%\n"
                result += f"  • Annualized Volatility: {annualized_volatility*100:.2f}%\n"
                result += f"  • Average Daily Return: {avg_return*100:+.2f}%\n"
                
                result += f"\nMomentum Indicators:\n"
                result += f"  • 7-day Trend: {recent_trend*100:+.2f}%\n"
                result += f"  • 3-day Trend: {short_term_trend*100:+.2f}%\n"
                result += f"  • 10-day Momentum: {momentum*100:+.2f}%\n"
                
                result += f"\n📊 Prediction Methods:\n"
                for method, pred_price in predictions:
                    change_pct = ((pred_price - current_price) / current_price) * 100
                    result += f"  • {method}: {format_price(pred_price)} ({change_pct:+.2f}%)"
                    if method == 'CryptoCraft.com' and cryptocraft_confidence:
                        result += f" [Confidence: {cryptocraft_confidence:.0f}%]"
                    result += "\n"
                
                if cryptocraft_prediction:
                    result += f"\n🌐 External Data Sources:\n"
                    result += f"  • CryptoCraft.com: Integrated prediction data\n"
            
            result += "\n⚠️  Risk Disclaimer:\n"
            result += "Price predictions are based on historical data and technical analysis.\n"
            result += "Past performance does not guarantee future results. Market conditions\n"
            result += "can change rapidly. Always do your own research and consider multiple\n"
            result += "factors before making investment decisions.\n"
            
            return [TextContent(type="text", text=result)]
            
        except ImportError as e:
            return [TextContent(
                type="text",
                text=f"Market data API not available: {str(e)}"
            )]
        except Exception as e:
            logger.error(f"Error predicting price: {e}", exc_info=True)
            return [TextContent(type="text", text=f"Error predicting price: {str(e)}")]
    
    async def _identify_image(self, image_data: str, query: Optional[str] = None) -> List[TextContent]:
        """
        Identify and describe the contents of an image using Ollama's vision API.
        
        This method analyzes images to detect:
        - Objects, people, scenes
        - Text content (OCR)
        - UI elements, buttons, icons
        - Code screenshots
        - Diagrams and charts
        - Screenshots and interface elements
        
        Args:
            image_data: Base64-encoded image data (data URL or raw base64). If empty, uses current images from context.
            query: Optional specific question about the image
            
        Returns:
            List of TextContent objects with image analysis results
        """
        try:
            import base64
            import json
            
            # If no image_data provided, try to use current images from context
            if not image_data or not image_data.strip():
                if self._current_images and len(self._current_images) > 0:
                    image_data = self._current_images[0]
                    logger.info(f"[MCP] Using current image from context (no image_data provided)")
                else:
                    return [TextContent(
                        type="text",
                        text="ERROR: No image data provided and no images are available in the current context. Please provide image_data or ensure images are attached to the message."
                    )]
            
            # Extract base64 data from data URL if needed
            if image_data.startswith("data:image"):
                # Format: data:image/png;base64,iVBORw0KGgo...
                base64_data = image_data.split(",", 1)[1] if "," in image_data else image_data
            else:
                base64_data = image_data
            
            # Validate base64 data
            try:
                # Try to decode to verify it's valid base64
                base64.b64decode(base64_data[:100])  # Just check first 100 chars
            except Exception as e:
                return [TextContent(
                    type="text",
                    text=f"ERROR: Invalid base64 image data: {str(e)}"
                )]
            
            # Get Ollama URL from environment or use default
            ollama_url = os.getenv("OLLAMA_URL", "http://localhost:5000")
            ollama_direct = os.getenv("OLLAMA_DIRECT_URL", "http://localhost:11434")
            
            # Try direct connection first, fallback to proxy
            ollama_base = ollama_direct
            
            # Build the prompt for image identification
            if query:
                prompt = f"{query}\n\nPlease analyze this image in detail and provide a comprehensive description."
            else:
                prompt = (
                    "Please analyze this image in detail and describe:\n"
                    "- What objects, people, or scenes are visible\n"
                    "- Any text content (if readable)\n"
                    "- UI elements, buttons, icons, or interface components (if it's a screenshot)\n"
                    "- Code or programming content (if visible)\n"
                    "- Diagrams, charts, or visualizations\n"
                    "- Overall context and purpose of the image\n"
                    "- Any other relevant details"
                )
            
            # Use a vision-capable model (try common vision models)
            vision_models = ["llava", "llava:latest", "llava:13b", "llava:7b", "bakllava", "llama3.2-vision"]
            model = None
            
            # Try to detect available vision model
            try:
                async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=10)) as session:
                    async with session.get(f"{ollama_base}/api/tags") as response:
                        if response.status == 200:
                            models_data = await response.json()
                            available_models = [m.get("name", "") for m in models_data.get("models", [])]
                            # Find first available vision model
                            for vm in vision_models:
                                if vm in available_models:
                                    model = vm
                                    break
            except Exception as e:
                logger.warning(f"Could not check available models: {e}")
            
            # Fallback to default vision model if none found
            if not model:
                model = vision_models[0]  # Default to llava
            
            # Call Ollama's vision API
            try:
                async with aiohttp.ClientSession(timeout=aiohttp.ClientTimeout(total=60)) as session:
                    payload = {
                        "model": model,
                        "prompt": prompt,
                        "images": [base64_data],
                        "stream": False
                    }
                    
                    async with session.post(
                        f"{ollama_base}/api/generate",
                        json=payload,
                        headers={"Content-Type": "application/json"}
                    ) as response:
                        if response.status == 200:
                            result_data = await response.json()
                            identification = result_data.get("response", "").strip()
                            
                            if identification:
                                result = f"Image Analysis Results:\n\n{identification}"
                                if query:
                                    result = f"Question: {query}\n\n{result}"
                                return [TextContent(type="text", text=result)]
                            else:
                                return [TextContent(
                                    type="text",
                                    text="Image analyzed but no description was returned. The image may be empty or the model may not have been able to process it."
                                )]
                        else:
                            error_text = await response.text()
                            return [TextContent(
                                type="text",
                                text=f"ERROR: Failed to analyze image. Ollama API returned status {response.status}: {error_text}"
                            )]
            except aiohttp.ClientError as e:
                return [TextContent(
                    type="text",
                    text=f"ERROR: Network error connecting to Ollama vision API: {str(e)}\n\nMake sure Ollama is running and a vision model (like llava) is installed."
                )]
            except asyncio.TimeoutError:
                return [TextContent(
                    type="text",
                    text="ERROR: Image analysis timed out. The vision model may be taking too long to process the image."
                )]
                
        except Exception as e:
            logger.error(f"Error identifying image: {e}", exc_info=True)
            return [TextContent(
                type="text",
                text=f"ERROR: Failed to identify image: {str(e)}"
            )]


def create_mcp_server(file_service=None, code_analyzer=None, web_search_enabled=True, location_service=None, memory_service=None):
    """
    Create and configure an MCP server instance following the Model Context Protocol
    
    This function creates a fully compliant MCP server that implements the
    standardized protocol for AI model integration. The server acts as a bridge
    between AI models and external systems, providing:
    
    - Unified Connectivity: Single interface for all tools and data sources
    - Real-Time Access: Live data retrieval instead of static training data
    - Scalability: Minimal setup, extensible architecture
    - Protocol Compliance: Follows MCP standard for maximum compatibility
    
    The server can be used in two modes:
    1. Direct integration (current): Tools accessed via MCPServerTools class
    2. Standalone server: Can run as separate process using stdio_server
    
    Example use case:
    Instead of writing separate integrations for Notion, Google Sheets, and
    project management tools, connect them all through this MCP server.
    The AI then queries the server, which fetches data and delivers it in real time.
    
    Args:
        file_service: Service for file operations
        code_analyzer: Service for code analysis
        web_search_enabled: Whether to enable web search tool
        location_service: Service for location detection, weather, and news
        
    Returns:
        Tuple of (Server instance, MCPServerTools instance) or (None, None) if MCP unavailable
    """
    if not MCP_AVAILABLE:
        return None
    
    tools = MCPServerTools(file_service, code_analyzer, web_search_enabled, location_service=location_service, memory_service=memory_service)
    server = Server("ai-agent-mcp")
    
    @server.list_tools()
    async def list_tools() -> List[Tool]:
        return tools.get_tools()
    
    @server.call_tool()
    async def call_tool(name: str, arguments: Dict[str, Any]) -> List[TextContent]:
        # Determine if write operations are allowed based on context
        # This will be set by the MCP client wrapper
        allow_write = arguments.pop("_allow_write", True)
        return await tools.execute_tool(name, arguments, allow_write=allow_write)
    
    return server, tools

